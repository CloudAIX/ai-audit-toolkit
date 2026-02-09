#!/usr/bin/env python3
"""
AI Audit Toolkit
Run professional AI audits for clients ($10K engagement framework).

Usage:
    python audit_toolkit.py --new-audit          # Start new audit project
    python audit_toolkit.py --questions          # Generate interview questions
    python audit_toolkit.py --roi                # Calculate ROI
    python audit_toolkit.py --report             # Generate audit report
"""

import argparse
import json
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass, asdict
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================================
# DATA MODELS
# ============================================================================

@dataclass
class Client:
    company_name: str
    industry: str
    employee_count: int
    contact_name: str
    contact_email: str
    avg_salary: float = 65000  # Annual average

@dataclass
class Opportunity:
    name: str
    description: str
    hours_saved_weekly: float
    employees_affected: int
    effort: str  # low, medium, high
    impact: str  # low, medium, high
    category: str = ""  # quick_win, big_swing, nice_to_have, deprioritize

    def __post_init__(self):
        # Auto-categorize based on effort/impact
        if not self.category:
            if self.effort == "low" and self.impact == "high":
                self.category = "quick_win"
            elif self.effort == "high" and self.impact == "high":
                self.category = "big_swing"
            elif self.effort == "low" and self.impact == "low":
                self.category = "nice_to_have"
            else:
                self.category = "deprioritize"

@dataclass
class AuditProject:
    client: Client
    opportunities: List[Opportunity]
    created_date: str
    interviews_completed: int = 0
    status: str = "discovery"  # discovery, analysis, presentation

# ============================================================================
# INTERVIEW QUESTION TEMPLATES
# ============================================================================

STAKEHOLDER_QUESTIONS = {
    "role_overview": [
        "Can you describe your role and your team's primary responsibilities?",
        "What are the main goals or KPIs your team is responsible for this quarter/year?",
        "Could you walk me through your team's structure?",
    ],
    "processes": [
        "From a high level, what are the most critical processes your team manages?",
        "Where do you see the biggest bottlenecks or delays in your team's workflow?",
        "Which tasks seem to consume the most man-hours or resources?",
    ],
    "technology": [
        "What are the main software systems or tools your team relies on?",
        "What are your biggest frustrations with your current technology stack?",
        "Are there important processes that happen outside of your main software?",
    ],
    "pain_points": [
        "What are the biggest challenges your team is facing right now?",
        "If you had a magic wand, what problem would you solve overnight?",
        "What is preventing your team from being more efficient?",
    ],
    "vision": [
        "Where do you see the biggest opportunities for improvement?",
        "How does your team generally respond to new technology?",
    ]
}

ENDUSER_QUESTIONS = {
    "daily_work": [
        "Can you walk me through a typical day or week in your role?",
        "What are the 1-3 most common tasks you perform every day?",
        "How much time is spent on core work versus administrative/repetitive tasks?",
    ],
    "process_deep_dive": [
        "Walk me through the exact steps to complete [specific common task]",
        "Which part is the most manual or takes the most time?",
        "What information do you need and where do you get it?",
    ],
    "tools": [
        "What software do you spend most of your day in?",
        "What do you find most frustrating about your tools?",
        "Is there double-entry or copying between systems?",
    ],
    "pain_points": [
        "What is the most boring or repetitive part of your job?",
        "If you had an assistant, what tasks would you give them immediately?",
        "How do you currently track and report on your work?",
    ]
}

INDUSTRY_SPECIFIC = {
    "healthcare": [
        "How do you currently handle patient intake and documentation?",
        "What compliance requirements create the most administrative burden?",
        "How much time is spent on insurance and billing-related tasks?",
        "What patient communication happens manually vs. automated?",
    ],
    "professional_services": [
        "How do you currently track billable hours and client work?",
        "What's your process for client onboarding?",
        "How much time goes into creating proposals and reports?",
        "What research or document review tasks are most time-consuming?",
    ],
    "retail_ecommerce": [
        "How do you handle inventory management and forecasting?",
        "What's your process for handling customer inquiries and returns?",
        "How do you currently manage product listings and pricing?",
        "What manual work goes into order fulfillment?",
    ],
    "finance": [
        "How do you handle data entry and reconciliation?",
        "What reporting tasks are most time-intensive?",
        "How do you currently manage compliance documentation?",
        "What client communication is done manually?",
    ],
    "manufacturing": [
        "How do you track production and quality metrics?",
        "What manual work goes into supply chain management?",
        "How do you handle equipment maintenance scheduling?",
        "What reporting and documentation is most time-consuming?",
    ],
    "aged_care": [
        "How do you currently manage incident reporting and SIRS notifications to the Aged Care Quality and Safety Commission?",
        "What's your process for medication management, round documentation, and controlled drug registers?",
        "How do you collect and organise evidence for the 8 Aged Care Quality Standards?",
        "What's your experience with AN-ACC funding model documentation and care minutes tracking?",
        "How do you manage clinical care plans, resident assessments, and progress notes?",
        "What systems do you use for staff rostering and ensuring you meet mandated care minute targets per resident?",
        "How do you communicate with residents' families about care updates, wellbeing, and incident notifications?",
        "What's your biggest compliance documentation challenge in preparing for Commission assessment contacts?",
        "How do you currently handle reportable assaults, unexplained absences, and other SIRS-reportable events?",
        "What proportion of your admin time goes to mandatory government reporting vs. direct care support?",
    ]
}

# ============================================================================
# ROI CALCULATOR
# ============================================================================

def calculate_roi(
    hours_saved_weekly: float,
    employees_affected: int,
    avg_annual_salary: float,
    implementation_cost: float,
    automation_efficiency: float = 0.7  # 70% of promised time savings realized
) -> dict:
    """Calculate ROI for an AI automation opportunity."""

    # Hourly rate
    hourly_rate = avg_annual_salary / 2080  # 40 hrs/week * 52 weeks

    # Adjusted hours saved (account for ramp-up and realistic efficiency)
    actual_hours_saved = hours_saved_weekly * employees_affected * automation_efficiency

    # Weekly and annual savings
    weekly_savings = actual_hours_saved * hourly_rate
    annual_savings = weekly_savings * 52

    # Revenue potential (50% of saved time redirected to revenue activities)
    revenue_hours_weekly = actual_hours_saved * 0.5
    # Assume revenue activities are worth 2x salary cost
    weekly_revenue_potential = revenue_hours_weekly * hourly_rate * 2
    annual_revenue_potential = weekly_revenue_potential * 52

    # Total annual value
    total_annual_value = annual_savings + annual_revenue_potential

    # ROI calculation
    roi_percentage = ((annual_savings - implementation_cost) / implementation_cost) * 100 if implementation_cost > 0 else 0
    payback_months = (implementation_cost / (annual_savings / 12)) if annual_savings > 0 else float('inf')

    return {
        "hourly_rate": round(hourly_rate, 2),
        "hours_saved_weekly": round(actual_hours_saved, 1),
        "weekly_savings": round(weekly_savings, 2),
        "annual_savings": round(annual_savings, 2),
        "annual_revenue_potential": round(annual_revenue_potential, 2),
        "total_annual_value": round(total_annual_value, 2),
        "roi_percentage": round(roi_percentage, 1),
        "payback_months": round(payback_months, 1),
        "implementation_cost": implementation_cost
    }

def calculate_audit_roi(opportunities: List[Opportunity], avg_salary: float, implementation_cost: float) -> dict:
    """Calculate combined ROI for all opportunities."""

    total_hours = sum(o.hours_saved_weekly * o.employees_affected for o in opportunities)
    total_employees = sum(o.employees_affected for o in opportunities)

    combined = calculate_roi(
        hours_saved_weekly=total_hours / max(total_employees, 1),
        employees_affected=total_employees,
        avg_annual_salary=avg_salary,
        implementation_cost=implementation_cost
    )

    # By category
    by_category = {}
    for cat in ["quick_win", "big_swing", "nice_to_have", "deprioritize"]:
        cat_opps = [o for o in opportunities if o.category == cat]
        if cat_opps:
            cat_hours = sum(o.hours_saved_weekly * o.employees_affected for o in cat_opps)
            cat_employees = sum(o.employees_affected for o in cat_opps)
            by_category[cat] = calculate_roi(
                hours_saved_weekly=cat_hours / max(cat_employees, 1),
                employees_affected=cat_employees,
                avg_annual_salary=avg_salary,
                implementation_cost=implementation_cost * 0.25  # Rough per-category estimate
            )

    return {
        "combined": combined,
        "by_category": by_category
    }

# ============================================================================
# REPORT GENERATOR
# ============================================================================

def generate_interview_doc(client: Client, role_type: str = "both") -> str:
    """Generate interview questions document for a client."""

    industry = client.industry.lower().replace(" ", "_")

    doc = f"""# AI Audit Interview Questions
## {client.company_name}

**Prepared for**: {client.contact_name}
**Date**: {datetime.now().strftime("%Y-%m-%d")}
**Industry**: {client.industry}
**Employee Count**: {client.employee_count}

---

## Interview Plan

| Business Size | Recommended Interviews | Duration |
|---------------|------------------------|----------|
| {client.employee_count} employees | {"3-5" if client.employee_count < 50 else "10-15"} interviews | 30-45 min each |

**Target mix**:
- 40% Leadership/Stakeholders (understand goals)
- 60% End-users (understand reality)

---

"""

    if role_type in ["both", "stakeholder"]:
        doc += "## Stakeholder Interview Questions (30,000-Foot View)\n\n"
        for section, questions in STAKEHOLDER_QUESTIONS.items():
            doc += f"### {section.replace('_', ' ').title()}\n\n"
            for q in questions:
                doc += f"- {q}\n"
            doc += "\n"

    if role_type in ["both", "enduser"]:
        doc += "## End-User Interview Questions (On-the-Ground Reality)\n\n"
        for section, questions in ENDUSER_QUESTIONS.items():
            doc += f"### {section.replace('_', ' ').title()}\n\n"
            for q in questions:
                doc += f"- {q}\n"
            doc += "\n"

    # Add industry-specific questions
    if industry in INDUSTRY_SPECIFIC:
        industry_display = client.industry.replace("_", " ").title()
        doc += f"## {industry_display} Industry-Specific Questions\n\n"
        for q in INDUSTRY_SPECIFIC[industry]:
            doc += f"- {q}\n"
        doc += "\n"

    doc += """---

## Interview Best Practices

- **Listen 80%, Talk 20%** — get them talking
- **Ask "Why?" repeatedly** — get to root causes
- **Record with permission** — use Fireflies.ai for transcription
- **Focus on problems, not solutions** — save solutions for later
- **Note emotional reactions** — frustration = opportunity

## After Each Interview

1. [ ] Transcription saved
2. [ ] Key pain points highlighted
3. [ ] Time estimates noted (hours/week on tasks)
4. [ ] Follow-up questions documented
"""

    return doc


def generate_opportunity_matrix(opportunities: List[Opportunity]) -> str:
    """Generate opportunity matrix visualization."""

    matrix = """# AI Opportunity Matrix

## Quick Reference

| Quadrant | Effort | Impact | Action |
|----------|--------|--------|--------|
| 🎯 Quick Wins | Low | High | **Priority #1** - Start here |
| 🚀 Big Swings | High | High | Long-term, high-ticket |
| ✨ Nice-to-Haves | Low | Low | Add-on value |
| ⏸️ Deprioritize | High | Low | Avoid |

---

## Identified Opportunities

"""

    categories = {
        "quick_win": ("🎯 Quick Wins (Low Effort, High Impact)", []),
        "big_swing": ("🚀 Big Swings (High Effort, High Impact)", []),
        "nice_to_have": ("✨ Nice-to-Haves (Low Effort, Low Impact)", []),
        "deprioritize": ("⏸️ Deprioritize (High Effort, Low Impact)", [])
    }

    for opp in opportunities:
        categories[opp.category][1].append(opp)

    for cat_id, (cat_name, opps) in categories.items():
        if opps:
            matrix += f"### {cat_name}\n\n"
            for opp in opps:
                matrix += f"**{opp.name}**\n"
                matrix += f"- {opp.description}\n"
                matrix += f"- Hours saved: {opp.hours_saved_weekly}/week × {opp.employees_affected} employees\n"
                matrix += f"- Effort: {opp.effort.upper()} | Impact: {opp.impact.upper()}\n\n"

    return matrix


def generate_executive_report(project: AuditProject, roi_data: dict) -> str:
    """Generate executive summary report."""

    client = project.client
    opportunities = project.opportunities

    quick_wins = [o for o in opportunities if o.category == "quick_win"]
    big_swings = [o for o in opportunities if o.category == "big_swing"]

    report = f"""# AI Audit Report
## {client.company_name}

**Prepared by**: GVRN-AI
**Date**: {datetime.now().strftime("%Y-%m-%d")}
**Engagement**: AI Opportunity Assessment

---

## Executive Summary

Following {project.interviews_completed} discovery interviews across {client.company_name},
we identified **{len(opportunities)} AI automation opportunities** with potential annual
value of **${roi_data['combined']['total_annual_value']:,.0f}**.

### Key Findings

| Metric | Value |
|--------|-------|
| Total Opportunities Identified | {len(opportunities)} |
| Quick Wins (Start Immediately) | {len(quick_wins)} |
| Strategic Initiatives | {len(big_swings)} |
| Estimated Hours Saved Weekly | {roi_data['combined']['hours_saved_weekly']:.0f} |
| Annual Cost Savings | ${roi_data['combined']['annual_savings']:,.0f} |
| Annual Revenue Potential | ${roi_data['combined']['annual_revenue_potential']:,.0f} |
| **Total Annual Value** | **${roi_data['combined']['total_annual_value']:,.0f}** |

---

## Recommended Roadmap

### Phase 1: Quick Wins (Weeks 1-4)

"""

    for i, opp in enumerate(quick_wins[:3], 1):
        report += f"""#### {i}. {opp.name}

- **Current State**: {opp.description}
- **Time Impact**: {opp.hours_saved_weekly} hours/week × {opp.employees_affected} people
- **Implementation**: 1-2 weeks

"""

    if big_swings:
        report += "### Phase 2: Strategic Initiatives (Months 2-6)\n\n"
        for i, opp in enumerate(big_swings[:3], 1):
            report += f"""#### {i}. {opp.name}

- **Current State**: {opp.description}
- **Time Impact**: {opp.hours_saved_weekly} hours/week × {opp.employees_affected} people
- **Implementation**: 4-8 weeks

"""

    report += f"""---

## ROI Analysis

### Cost Savings Calculation

```
Hours Saved/Week:     {roi_data['combined']['hours_saved_weekly']:.0f} hours
Average Hourly Rate:  ${roi_data['combined']['hourly_rate']:.2f}
Weekly Savings:       ${roi_data['combined']['weekly_savings']:,.0f}
Annual Savings:       ${roi_data['combined']['annual_savings']:,.0f}
```

### Revenue Potential

Assuming 50% of saved time is redirected to revenue-generating activities:

```
Annual Revenue Potential: ${roi_data['combined']['annual_revenue_potential']:,.0f}
```

### Investment Summary

| Metric | Value |
|--------|-------|
| Estimated Implementation Cost | ${roi_data['combined']['implementation_cost']:,.0f} |
| Payback Period | {roi_data['combined']['payback_months']:.1f} months |
| First Year ROI | {roi_data['combined']['roi_percentage']:.0f}% |

---

## Next Steps

1. **Approve Phase 1 Quick Wins** — Start with highest-impact, lowest-effort items
2. **Schedule kickoff meeting** — Align team and set success metrics
3. **Begin implementation** — Target 2-4 week delivery for first automation

---

*Report generated by GVRN-AI | AI Audit Framework*
*Contact: [your-email@gvrn-ai.com]*
"""

    return report


# ============================================================================
# POWERPOINT REPORT GENERATOR
# ============================================================================

# Brand colours
_DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
_GREEN_ACC   = RGBColor(0x00, 0xC9, 0x7B)   # GVRN green
_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GREY  = RGBColor(0xCC, 0xCC, 0xCC)
_MID_GREY    = RGBColor(0x88, 0x88, 0x88)
_TABLE_ROW1  = RGBColor(0x22, 0x22, 0x3A)   # table odd-row
_TABLE_ROW2  = RGBColor(0x2A, 0x2A, 0x44)   # table even-row

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, colour=_DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_textbox(slide, left, top, width, height, text,
                 font_size=18, colour=_WHITE, bold=False,
                 alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = colour
    p.font.bold = bold
    p.alignment = alignment
    return tf


def _add_table(slide, rows_data, left, top, width, col_widths,
               header_colour=_GREEN_ACC, row_font_size=14):
    """Add a styled table to a slide. rows_data = list of tuples."""
    rows = len(rows_data)
    cols = len(rows_data[0]) if rows_data else 2
    table_shape = slide.shapes.add_table(rows, cols, left, top, width,
                                          Emu(rows * Pt(row_font_size + 18).emu))
    table = table_shape.table

    # Set column widths
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw

    for r_idx, row_vals in enumerate(rows_data):
        for c_idx, val in enumerate(row_vals):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            # Cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if r_idx == 0:
                cell_fill.fore_color.rgb = _DARK_BG
            else:
                cell_fill.fore_color.rgb = _TABLE_ROW1 if r_idx % 2 == 1 else _TABLE_ROW2

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(row_font_size)
                paragraph.font.color.rgb = header_colour if r_idx == 0 else _WHITE
                paragraph.font.bold = r_idx == 0
                paragraph.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.RIGHT

    # Remove table borders for a cleaner look via XML
    from pptx.oxml.ns import qn
    for r_idx in range(rows):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                ln = tcPr.find(qn(border_tag))
                if ln is not None:
                    tcPr.remove(ln)
                from lxml import etree
                ln_el = etree.SubElement(tcPr, qn(border_tag), w="0", cap="flat")
                etree.SubElement(ln_el, qn("a:noFill"))

    return table


def _slide_title_bar(slide, title_text, subtitle_text=None):
    """Add a thin green accent line + title at the top of a content slide."""
    # Green accent line
    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(0.5),
                                   Inches(0.08), Inches(0.55))
    line.fill.solid()
    line.fill.fore_color.rgb = _GREEN_ACC
    line.line.fill.background()

    _add_textbox(slide, Inches(0.9), Inches(0.4), Inches(10), Inches(0.55),
                 title_text, font_size=28, colour=_WHITE, bold=True)
    if subtitle_text:
        _add_textbox(slide, Inches(0.9), Inches(0.95), Inches(10), Inches(0.4),
                     subtitle_text, font_size=14, colour=_LIGHT_GREY)


def generate_executive_pptx(project: AuditProject, roi_data: dict,
                            output_path: Path) -> Path:
    """Generate a branded executive PowerPoint presentation."""

    client = project.client
    opportunities = project.opportunities
    quick_wins = [o for o in opportunities if o.category == "quick_win"]
    big_swings = [o for o in opportunities if o.category == "big_swing"]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ------------------------------------------------------------------
    # SLIDE 1: Cover
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)

    # Green accent bar at top
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0),
                                  SLIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _GREEN_ACC
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.9),
                 "GVRN-AI", font_size=20, colour=_GREEN_ACC, bold=True)
    _add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.0),
                 "AI Opportunity Assessment", font_size=40, colour=_WHITE, bold=True)
    _add_textbox(slide, Inches(0.8), Inches(3.6), Inches(11), Inches(0.7),
                 client.company_name, font_size=28, colour=_LIGHT_GREY)

    date_str = datetime.now().strftime("%B %Y")
    _add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
                 f"Prepared for {client.contact_name}  |  {date_str}",
                 font_size=14, colour=_MID_GREY)

    # Bottom green bar
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), SLIDE_H - Inches(0.08),
                                   SLIDE_W, Inches(0.08))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = _GREEN_ACC
    bar2.line.fill.background()

    # ------------------------------------------------------------------
    # SLIDE 2: Executive Summary
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _slide_title_bar(slide, "Executive Summary",
                     f"Based on {project.interviews_completed} discovery interviews")

    combined = roi_data["combined"]
    summary_rows = [
        ("Metric", "Value"),
        ("Total Opportunities Identified", str(len(opportunities))),
        ("Quick Wins (Start Immediately)", str(len(quick_wins))),
        ("Strategic Initiatives", str(len(big_swings))),
        ("Estimated Hours Saved / Week", f"{combined['hours_saved_weekly']:.0f} hrs"),
        ("Annual Cost Savings", f"${combined['annual_savings']:,.0f}"),
        ("Annual Revenue Potential", f"${combined['annual_revenue_potential']:,.0f}"),
        ("Total Annual Value", f"${combined['total_annual_value']:,.0f}"),
    ]

    _add_table(slide, summary_rows,
               left=Inches(0.8), top=Inches(1.6),
               width=Inches(8), col_widths=[Inches(5), Inches(3)],
               row_font_size=15)

    # Callout box on the right
    _add_textbox(slide, Inches(9.3), Inches(2.0), Inches(3.5), Inches(0.4),
                 "TOTAL ANNUAL VALUE", font_size=13, colour=_GREEN_ACC, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(9.3), Inches(2.5), Inches(3.5), Inches(0.7),
                 f"${combined['total_annual_value']:,.0f}",
                 font_size=36, colour=_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(9.3), Inches(3.3), Inches(3.5), Inches(0.4),
                 f"ROI {combined['roi_percentage']:.0f}%  |  Payback {combined['payback_months']:.1f} mo",
                 font_size=13, colour=_LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # SLIDES 3–N: One per Quick Win
    # ------------------------------------------------------------------
    for idx, opp in enumerate(quick_wins, 1):
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)
        _slide_title_bar(slide, f"Quick Win #{idx}: {opp.name}",
                         "Low Effort  |  High Impact")

        # Current State
        _add_textbox(slide, Inches(0.9), Inches(1.6), Inches(5.5), Inches(0.35),
                     "CURRENT STATE", font_size=13, colour=_GREEN_ACC, bold=True)
        _add_textbox(slide, Inches(0.9), Inches(2.0), Inches(5.5), Inches(1.5),
                     opp.description, font_size=16, colour=_LIGHT_GREY)

        # Proposed Solution
        _add_textbox(slide, Inches(0.9), Inches(3.7), Inches(5.5), Inches(0.35),
                     "PROPOSED SOLUTION", font_size=13, colour=_GREEN_ACC, bold=True)

        solution_text = (f"Automate and streamline this workflow using AI-powered tools, "
                         f"eliminating manual effort and reducing error rates.")
        _add_textbox(slide, Inches(0.9), Inches(4.1), Inches(5.5), Inches(1.2),
                     solution_text, font_size=16, colour=_LIGHT_GREY)

        # Impact panel on the right
        hourly = roi_data["combined"]["hourly_rate"]
        weekly_hrs = opp.hours_saved_weekly * opp.employees_affected * 0.7
        weekly_saving = weekly_hrs * hourly
        annual_saving = weekly_saving * 52

        impact_rows = [
            ("Impact Metric", "Value"),
            ("Hours Saved / Week", f"{weekly_hrs:.0f} hrs"),
            ("Employees Affected", str(opp.employees_affected)),
            ("Est. Weekly Saving", f"${weekly_saving:,.0f}"),
            ("Est. Annual Saving", f"${annual_saving:,.0f}"),
            ("Implementation", "1-2 weeks"),
        ]
        _add_table(slide, impact_rows,
                   left=Inches(7.0), top=Inches(1.6),
                   width=Inches(5.5), col_widths=[Inches(3), Inches(2.5)],
                   row_font_size=14)

    # ------------------------------------------------------------------
    # SLIDE: ROI Summary
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _slide_title_bar(slide, "Return on Investment",
                     "Cost savings + revenue potential from redirected capacity")

    roi_rows = [
        ("Metric", "Value"),
        ("Estimated Implementation Cost", f"${combined['implementation_cost']:,.0f}"),
        ("Hours Saved / Week", f"{combined['hours_saved_weekly']:.0f} hrs"),
        ("Annual Cost Savings", f"${combined['annual_savings']:,.0f}"),
        ("Annual Revenue Potential", f"${combined['annual_revenue_potential']:,.0f}"),
        ("Total Annual Value", f"${combined['total_annual_value']:,.0f}"),
        ("Payback Period", f"{combined['payback_months']:.1f} months"),
        ("First Year ROI", f"{combined['roi_percentage']:.0f}%"),
    ]
    _add_table(slide, roi_rows,
               left=Inches(0.8), top=Inches(1.6),
               width=Inches(8), col_widths=[Inches(5), Inches(3)],
               row_font_size=15)

    # Visual callout
    _add_textbox(slide, Inches(9.3), Inches(2.0), Inches(3.5), Inches(0.4),
                 "FIRST YEAR ROI", font_size=13, colour=_GREEN_ACC, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(9.3), Inches(2.5), Inches(3.5), Inches(0.7),
                 f"{combined['roi_percentage']:.0f}%",
                 font_size=44, colour=_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(9.3), Inches(3.4), Inches(3.5), Inches(0.4),
                 f"Payback in {combined['payback_months']:.1f} months",
                 font_size=14, colour=_LIGHT_GREY, alignment=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # SLIDE: Next Steps
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide)
    _slide_title_bar(slide, "Recommended Next Steps")

    steps = [
        ("1", "Approve Phase 1 Quick Wins",
         "Start with the highest-impact, lowest-effort items identified in this assessment."),
        ("2", "Schedule Kickoff Meeting",
         "Align the project team and define success metrics for each automation."),
        ("3", "Begin Implementation",
         "Target 2-4 week delivery for the first automation, with iterative rollout."),
        ("4", "Measure & Iterate",
         "Track time savings against baseline and expand to Phase 2 strategic initiatives."),
    ]

    y_offset = Inches(1.8)
    for num, title, desc in steps:
        # Number circle (green square as proxy)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(0.9), y_offset,
                                         Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = _GREEN_ACC
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.color.rgb = _DARK_BG
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

        _add_textbox(slide, Inches(1.6), y_offset - Inches(0.05),
                     Inches(10), Inches(0.4),
                     title, font_size=20, colour=_WHITE, bold=True)
        _add_textbox(slide, Inches(1.6), y_offset + Inches(0.35),
                     Inches(10), Inches(0.4),
                     desc, font_size=14, colour=_LIGHT_GREY)

        y_offset += Inches(1.2)

    # Footer on last slide
    _add_textbox(slide, Inches(0.8), SLIDE_H - Inches(0.7), Inches(11), Inches(0.4),
                 "GVRN-AI  |  AI Audit Framework  |  Confidential",
                 font_size=11, colour=_MID_GREY, alignment=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------
    prs.save(str(output_path))
    return output_path


# ============================================================================
# CLI INTERFACE
# ============================================================================

def interactive_new_audit() -> AuditProject:
    """Interactive wizard to create new audit project."""

    print("\n" + "="*60)
    print("  NEW AI AUDIT PROJECT")
    print("="*60 + "\n")

    # Client info
    company_name = input("Company name: ").strip()
    industry = input("Industry (healthcare, professional_services, retail_ecommerce, finance, manufacturing): ").strip()
    employee_count = int(input("Employee count: ").strip() or "50")
    contact_name = input("Primary contact name: ").strip()
    contact_email = input("Contact email: ").strip()
    avg_salary = float(input("Average annual salary [$65000]: ").strip() or "65000")

    client = Client(
        company_name=company_name,
        industry=industry,
        employee_count=employee_count,
        contact_name=contact_name,
        contact_email=contact_email,
        avg_salary=avg_salary
    )

    project = AuditProject(
        client=client,
        opportunities=[],
        created_date=datetime.now().isoformat()
    )

    return project


def interactive_add_opportunity() -> Opportunity:
    """Add an opportunity interactively."""

    print("\n--- Add Opportunity ---\n")

    name = input("Opportunity name: ").strip()
    description = input("Description (current problem): ").strip()
    hours_saved = float(input("Hours saved per week (per person): ").strip() or "5")
    employees = int(input("Number of employees affected: ").strip() or "1")
    effort = input("Effort (low/medium/high): ").strip().lower() or "medium"
    impact = input("Impact (low/medium/high): ").strip().lower() or "medium"

    return Opportunity(
        name=name,
        description=description,
        hours_saved_weekly=hours_saved,
        employees_affected=employees,
        effort=effort,
        impact=impact
    )


def save_project(project: AuditProject, output_dir: Path):
    """Save project to JSON."""

    output_dir.mkdir(parents=True, exist_ok=True)

    safe_name = project.client.company_name.lower().replace(" ", "_")
    filepath = output_dir / f"{safe_name}_audit.json"

    with open(filepath, "w") as f:
        json.dump(asdict(project), f, indent=2)

    print(f"Project saved to: {filepath}")
    return filepath


def load_project(filepath: Path) -> AuditProject:
    """Load project from JSON."""

    with open(filepath) as f:
        data = json.load(f)

    client = Client(**data["client"])
    opportunities = [Opportunity(**o) for o in data["opportunities"]]

    return AuditProject(
        client=client,
        opportunities=opportunities,
        created_date=data["created_date"],
        interviews_completed=data.get("interviews_completed", 0),
        status=data.get("status", "discovery")
    )


def main():
    parser = argparse.ArgumentParser(description="AI Audit Toolkit")
    parser.add_argument("--new-audit", action="store_true", help="Start new audit project")
    parser.add_argument("--questions", action="store_true", help="Generate interview questions")
    parser.add_argument("--roi", action="store_true", help="Calculate ROI")
    parser.add_argument("--report", action="store_true", help="Generate full report")
    parser.add_argument("--project", "-p", type=str, help="Path to project JSON file")
    parser.add_argument("--output", "-o", type=str, help="Output directory")
    parser.add_argument("--example", "-e", action="store_true", help="Run with example data")

    args = parser.parse_args()

    output_dir = Path(args.output) if args.output else Path.cwd() / "output"

    # Example mode
    if args.example:
        client = Client(
            company_name="Maplewood Residential Aged Care",
            industry="aged_care",
            employee_count=40,
            contact_name="Karen Mitchell",
            contact_email="karen.mitchell@maplewoodcare.com.au",
            avg_salary=62000  # AUD average across RNs, ENs, PCAs, and admin
        )

        opportunities = [
            Opportunity(
                "Digital Incident Reporting & SIRS Compliance",
                "Paper-based incident forms take 30-45 min each; SIRS notifications to the Aged Care Quality and Safety Commission are manually tracked in a spreadsheet",
                6, 3, "low", "high"
            ),
            Opportunity(
                "AI-Assisted Quality Standards Documentation",
                "Admin staff spend 12+ hrs/week manually compiling evidence portfolios for the 8 Aged Care Quality Standards and continuous improvement registers",
                10, 3, "low", "high"
            ),
            Opportunity(
                "Automated AN-ACC Care Minutes Tracking",
                "Manual tracking of direct and indirect care minutes across shifts for AN-ACC funding submissions; staff record on paper timesheets then admin re-enters into government portal",
                5, 5, "low", "high"
            ),
            Opportunity(
                "Electronic Medication Management",
                "Paper medication charts with manual round tracking; double-handling between pharmacy orders, GP scripts, and MAR charts increases medication error risk",
                4, 8, "high", "high"
            ),
            Opportunity(
                "Clinical Care Plan Automation",
                "Quarterly care plan reviews done manually across 45 residents with paper-based assessments; RNs spend evenings updating plans instead of providing direct care",
                3, 8, "high", "high"
            ),
            Opportunity(
                "Resident & Family Communication Portal",
                "Manual phone calls and printed letters to families for care updates, activity schedules, and incident notifications; families frequently call reception for updates",
                3, 4, "low", "low"
            ),
            Opportunity(
                "Staff Rostering Optimisation",
                "Manual roster creation in spreadsheets; difficulty balancing AN-ACC care minute targets, staff availability, and award conditions across 24/7 shifts",
                5, 2, "low", "low"
            ),
        ]

        project = AuditProject(
            client=client,
            opportunities=opportunities,
            created_date=datetime.now().isoformat(),
            interviews_completed=8,
            status="analysis"
        )

        # Generate all outputs
        output_dir.mkdir(parents=True, exist_ok=True)

        # Interview questions
        questions = generate_interview_doc(client)
        q_path = output_dir / "interview_questions.md"
        q_path.write_text(questions)
        print(f"Interview questions: {q_path}")

        # Opportunity matrix
        matrix = generate_opportunity_matrix(opportunities)
        m_path = output_dir / "opportunity_matrix.md"
        m_path.write_text(matrix)
        print(f"Opportunity matrix: {m_path}")

        # ROI calculation
        roi_data = calculate_audit_roi(opportunities, client.avg_salary, 25000)

        # Executive report
        report = generate_executive_report(project, roi_data)
        r_path = output_dir / "executive_report.md"
        r_path.write_text(report)
        print(f"Executive report: {r_path}")

        # Executive PowerPoint
        pptx_path = output_dir / "executive_presentation.pptx"
        generate_executive_pptx(project, roi_data, pptx_path)
        print(f"Executive PPTX:   {pptx_path}")

        # Save project
        save_project(project, output_dir)

        print("\n" + "="*60)
        print("Example audit generated! Check the output/ folder.")
        print("="*60)
        return

    # New audit
    if args.new_audit:
        project = interactive_new_audit()

        # Ask if they want to add opportunities
        while True:
            add_more = input("\nAdd opportunity? (y/n): ").strip().lower()
            if add_more == "y":
                opp = interactive_add_opportunity()
                project.opportunities.append(opp)
            else:
                break

        save_project(project, output_dir)
        return

    # Load existing project for other commands
    if args.project:
        project = load_project(Path(args.project))
    else:
        print("Use --example for demo, --new-audit to start, or --project <file> to load existing.")
        return

    # Generate questions
    if args.questions:
        questions = generate_interview_doc(project.client)
        q_path = output_dir / "interview_questions.md"
        q_path.write_text(questions)
        print(f"Saved to: {q_path}")

    # Calculate ROI
    if args.roi:
        roi_data = calculate_audit_roi(project.opportunities, project.client.avg_salary, 15000)
        print(json.dumps(roi_data, indent=2))

    # Generate report
    if args.report:
        roi_data = calculate_audit_roi(project.opportunities, project.client.avg_salary, 15000)
        report = generate_executive_report(project, roi_data)
        r_path = output_dir / "executive_report.md"
        r_path.write_text(report)
        print(f"Saved to: {r_path}")

        pptx_path = output_dir / "executive_presentation.pptx"
        generate_executive_pptx(project, roi_data, pptx_path)
        print(f"Saved to: {pptx_path}")


if __name__ == "__main__":
    main()
